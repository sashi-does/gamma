"""
PPT Generator MCP Server
Provides PowerPoint creation tools via FastMCP.
Maintains an internal Presentation object across tool calls.
"""

import os
import logging
from pathlib import Path
from typing import List

from fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Logging ────────────────────────────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="[PPT-MCP] %(asctime)s %(levelname)s │ %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger("ppt_mcp")

# ── Server ─────────────────────────────────────────────────────────────────────
mcp = FastMCP(name="ppt-generator-mcp")

# ── Internal state ─────────────────────────────────────────────────────────────
_state: dict = {
    "presentation": None,
    "filename": None,
    "slide_count": 0,
}

# ── Colour Palette — "Aurora Executive" ───────────────────────────────────────
BG_DARK       = RGBColor(0x0D, 0x1B, 0x2A)   # #0D1B2A  deep midnight
BG_PANEL      = RGBColor(0x11, 0x25, 0x40)   # #112540  dark navy panel
BG_LIGHT      = RGBColor(0xF5, 0xF7, 0xFA)   # #F5F7FA  near-white
ACCENT1       = RGBColor(0x00, 0xC8, 0xAA)   # #00C8AA  teal/mint
ACCENT2       = RGBColor(0x7B, 0x5E, 0xF8)   # #7B5EF8  violet
ACCENT3       = RGBColor(0xFF, 0x6B, 0x6B)   # #FF6B6B  coral (used sparingly)
HEADER_CLR    = RGBColor(0x1A, 0x2E, 0x4A)   # #1A2E4A  dark header bg
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BODY_TEXT     = RGBColor(0x1A, 0x2E, 0x4A)   # dark navy text on light bg
SUBTEXT       = RGBColor(0x55, 0x70, 0x90)   # muted slate for secondary text
CARD_COLORS   = [
    RGBColor(0x00, 0xC8, 0xAA),   # teal
    RGBColor(0x7B, 0x5E, 0xF8),   # violet
    RGBColor(0xFF, 0x6B, 0x6B),   # coral
    RGBColor(0xF7, 0xB7, 0x31),   # amber
    RGBColor(0x3A, 0xC4, 0xFF),   # sky blue
]

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────────────

def _solid_bg(slide, colour: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _add_rect(slide, left, top, width, height, fill_color: RGBColor, no_line=True):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if no_line:
        shape.line.fill.background()
    return shape


def _add_oval(slide, left, top, width, height, fill_color: RGBColor):
    shape = slide.shapes.add_shape(9, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_text_box(
    slide,
    left: float, top: float, width: float, height: float,
    text: str,
    font_size: int,
    bold: bool,
    colour: RGBColor,
    align=PP_ALIGN.LEFT,
    italic: bool = False,
    font_name: str = "Calibri",
) -> None:
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    run.font.name = font_name


# ── Title Slide ────────────────────────────────────────────────────────────────

def _build_title_slide(prs: Presentation, title: str) -> None:
    """
    Visually rich title slide:
    - Deep midnight background
    - Large decorative circle (top-right, semi-transparent effect via layering)
    - Teal horizontal band (bottom third)
    - Violet accent bar (left)
    - Bold title + subtitle label
    - Decorative dots bottom-left
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, BG_DARK)

    # ── Background geometry ─────────────────────────────────────────────────

    # Large decorative circle — top-right corner (gives depth)
    _add_oval(slide, left=9.5, top=-1.5, width=5.5, height=5.5, fill_color=BG_PANEL)

    # Smaller circle layered on top for glow illusion
    _add_oval(slide, left=10.5, top=-0.8, width=3.5, height=3.5, fill_color=RGBColor(0x16, 0x32, 0x55))

    # Teal wave-like bottom band (wide flat rectangle)
    _add_rect(slide, left=0, top=5.9, width=13.33, height=1.6, fill_color=RGBColor(0x00, 0x8A, 0x72))

    # Slightly lighter teal strip on top of band
    _add_rect(slide, left=0, top=5.9, width=13.33, height=0.18, fill_color=ACCENT1)

    # Violet accent bar — left edge
    _add_rect(slide, left=0, top=0, width=0.18, height=7.5, fill_color=ACCENT2)

    # Decorative small squares — bottom-left corner
    for i, color in enumerate([ACCENT1, ACCENT2, ACCENT3]):
        _add_rect(slide, left=0.35 + i * 0.32, top=6.9, width=0.22, height=0.22, fill_color=color)

    # Horizontal rule above title
    _add_rect(slide, left=0.4, top=2.35, width=4.5, height=0.04, fill_color=ACCENT1)

    # ── Text ────────────────────────────────────────────────────────────────

    # Main title
    _add_text_box(
        slide, left=0.4, top=2.5, width=11.5, height=1.8,
        text=title, font_size=48, bold=True,
        colour=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri Light",
    )

    # Accent title underline dots
    for i in range(3):
        _add_oval(slide, left=0.4 + i * 0.28, top=4.35, width=0.12, height=0.12, fill_color=ACCENT1)

    # Sub-label
    _add_text_box(
        slide, left=0.4, top=4.6, width=9, height=0.5,
        text="Auto-Generated by PPT Agent",
        font_size=15, bold=False, italic=True,
        colour=ACCENT1, align=PP_ALIGN.LEFT,
    )

    # Bottom-band watermark text
    _add_text_box(
        slide, left=0.4, top=6.05, width=10, height=0.4,
        text="CONFIDENTIAL  •  INTERNAL USE ONLY",
        font_size=10, bold=False, italic=False,
        colour=RGBColor(0xCC, 0xFF, 0xF0), align=PP_ALIGN.LEFT,
    )

    log.info("Built title slide: %s", title)


# ── Content Slide ──────────────────────────────────────────────────────────────

def _build_content_slide(
    prs: Presentation, title: str, bullets: List[str], slide_num: int
) -> None:
    """
    Rich content slide:
    - Light background with subtle dark side panel
    - Gradient-style layered header (dark + teal strip)
    - Card-style numbered bullet rows with colour-coded left tabs
    - Decorative background circle (bottom-right)
    - Footer with slide number and teal line
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, BG_LIGHT)

    # ── Background decorations ───────────────────────────────────────────────

    # Dark right panel (gives a two-tone layout feel)
    _add_rect(slide, left=10.5, top=1.4, width=2.83, height=5.75, fill_color=RGBColor(0xE8, 0xEC, 0xF2))

    # Large faint circle — bottom-right (decorative depth)
    _add_oval(slide, left=9.8, top=4.2, width=4.0, height=4.0, fill_color=RGBColor(0xDC, 0xE3, 0xED))

    # ── Header band ─────────────────────────────────────────────────────────

    # Main dark header
    _add_rect(slide, left=0, top=0, width=13.33, height=1.35, fill_color=HEADER_CLR)

    # Teal accent strip at the very bottom of header
    _add_rect(slide, left=0, top=1.28, width=13.33, height=0.1, fill_color=ACCENT1)

    # Violet left bar in header
    _add_rect(slide, left=0, top=0, width=0.18, height=1.35, fill_color=ACCENT2)

    # Slide number circle in header (top-right)
    _add_oval(slide, left=12.55, top=0.18, width=0.6, height=0.6, fill_color=ACCENT1)
    _add_text_box(
        slide, left=12.55, top=0.17, width=0.6, height=0.6,
        text=str(slide_num), font_size=14, bold=True,
        colour=BG_DARK, align=PP_ALIGN.CENTER,
    )

    # Title text in header
    _add_text_box(
        slide, left=0.35, top=0.18, width=11.8, height=1.0,
        text=title, font_size=30, bold=True,
        colour=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri Light",
    )

    # ── Card-style bullet rows ───────────────────────────────────────────────

    card_top_start = 1.55
    card_height    = 0.82
    card_gap       = 0.14
    card_width     = 10.0   # leaves right panel visible

    for i, bp in enumerate(bullets[:5]):
        color     = CARD_COLORS[i % len(CARD_COLORS)]
        card_top  = card_top_start + i * (card_height + card_gap)

        # Card background (white with slight shadow illusion via offset rect)
        _add_rect(slide, left=0.42, top=card_top + 0.04, width=card_width, height=card_height,
                  fill_color=RGBColor(0xD8, 0xDE, 0xE8))   # shadow
        _add_rect(slide, left=0.38, top=card_top, width=card_width, height=card_height,
                  fill_color=WHITE)

        # Coloured left tab on card
        _add_rect(slide, left=0.38, top=card_top, width=0.22, height=card_height, fill_color=color)

        # Number badge
        _add_oval(slide, left=0.72, top=card_top + 0.18, width=0.44, height=0.44, fill_color=color)
        _add_text_box(
            slide, left=0.72, top=card_top + 0.16, width=0.44, height=0.44,
            text=str(i + 1), font_size=13, bold=True,
            colour=WHITE, align=PP_ALIGN.CENTER,
        )

        # Bullet text
        _add_text_box(
            slide,
            left=1.28, top=card_top + 0.1,
            width=card_width - 1.0, height=card_height - 0.12,
            text=bp, font_size=17, bold=False,
            colour=BODY_TEXT, align=PP_ALIGN.LEFT,
        )

    # ── Footer bar ───────────────────────────────────────────────────────────

    _add_rect(slide, left=0, top=7.22, width=13.33, height=0.28, fill_color=HEADER_CLR)
    _add_rect(slide, left=0, top=7.22, width=13.33, height=0.05, fill_color=ACCENT1)

    _add_text_box(
        slide, left=0.3, top=7.23, width=8, height=0.25,
        text="AUTO-GENERATED PRESENTATION  •  CONFIDENTIAL",
        font_size=8, bold=False, italic=False,
        colour=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.LEFT,
    )

    log.info("Built content slide #%d: %s (%d bullets)", slide_num, title, len(bullets))


# ── MCP Tools ──────────────────────────────────────────────────────────────────

@mcp.tool()
def create_presentation(filename: str) -> str:
    """
    Initialise a new PowerPoint presentation with *filename* as the save path.
    Resets any previously open presentation.
    """
    global _state
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _state["presentation"] = prs
    _state["filename"]     = filename
    _state["slide_count"]  = 0

    log.info("create_presentation → %s", filename)
    return f"Presentation initialised. Will save to: {filename}"


@mcp.tool()
def add_slide(title: str, bullet_points: List[str]) -> str:
    """
    Add one slide to the current presentation.

    Args:
        title:         Slide heading (string).
        bullet_points: List of 3-5 bullet point strings.

    Returns a confirmation string.
    """
    prs = _state.get("presentation")
    if prs is None:
        return "ERROR: No presentation open. Call create_presentation first."

    if len(bullet_points) < 3:
        while len(bullet_points) < 3:
            bullet_points.append("Additional detail to be expanded upon.")
    bullet_points = bullet_points[:5]

    _state["slide_count"] += 1
    count = _state["slide_count"]

    if count == 1:
        _build_title_slide(prs, title)
    else:
        _build_content_slide(prs, title, bullet_points, count - 1)

    return (
        f"Slide {count} added: '{title}' "
        f"with {len(bullet_points)} bullet(s). "
        f"Total slides so far: {count}."
    )


@mcp.tool()
def save_presentation() -> str:
    """
    Save the current presentation to the filename given in create_presentation.
    Returns the full path of the saved file.
    """
    prs      = _state.get("presentation")
    filename = _state.get("filename")

    if prs is None:
        return "ERROR: No presentation open."
    if not filename:
        return "ERROR: No filename set. Call create_presentation first."

    save_path = Path(os.getcwd()) / filename
    prs.save(str(save_path))
    log.info("save_presentation → %s  (%d slides)", save_path, _state["slide_count"])
    return f"Presentation saved: {save_path}  ({_state['slide_count']} slides)"


# ── Entry point ────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    log.info("Starting PPT Generator MCP Server …")
    mcp.run(transport="stdio")